import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
import streamlit as st

def generate_sample_chart(data: dict, title="분석 그래프") -> str:
    """
    data: {'x': [...], 'y': [...]}
    간단한 시계열 선 그래프를 생성하여 이미지 파일(chart.png)로 저장 후 경로를 반환합니다.
    """
    try:
        plt.figure(figsize=(8, 4))
        plt.plot(data["x"], data["y"], marker="o", linestyle="-")
        plt.title(title)
        plt.xlabel("시간")
        plt.ylabel("값")
        plt.grid(True)
        chart_path = "chart.png"
        plt.savefig(chart_path, bbox_inches="tight")
        plt.close()
        return chart_path
    except Exception as e:
        st.error(f"차트 생성 실패: {e}")
        return None

def create_ppt_report(analysis_text: str, sql_code: str, chart_image_path: str, output_filename="OCI_분석보고서.pptx") -> str:
    """
    분석 인사이트, PostgreSQL 코드, 그리고 시각화 이미지를 포함한 PPT 보고서를 생성합니다.
    파일은 output_filename으로 저장되며, 생성된 파일 경로를 반환합니다.
    """
    prs = Presentation()

    # 1. 제목 슬라이드 (레이아웃 0)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = "데이터 분석 보고서"
    else:
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.5))
        title_box.text_frame.text = "데이터 분석 보고서"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "AI 기반 PostgreSQL 코드 및 분석 인사이트"

    # 2. 코드 및 인사이트 슬라이드 (레이아웃 5)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_found = False
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip() != "":
            shape.text = "PostgreSQL 코드 및 분석 인사이트"
            title_found = True
            break
    if not title_found:
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.5))
        textbox.text_frame.text = "PostgreSQL 코드 및 분석 인사이트"
    
    contentBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(3))
    contentFrame = contentBox.text_frame
    contentFrame.text = f"PostgreSQL 코드:\n{sql_code}\n\n분석 인사이트:\n{analysis_text}"
    for paragraph in contentFrame.paragraphs:
        paragraph.font.size = Pt(14)

    # 3. 시각화 슬라이드 (레이아웃 6)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_set = False
    for shape in slide.shapes:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text.strip() != "":
            shape.text = "분석 결과 시각화"
            title_set = True
            break
    if not title_set:
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.5))
        textbox.text_frame.text = "분석 결과 시각화"
    
    slide.shapes.add_picture(chart_image_path, Inches(1), Inches(2), width=Inches(8), height=Inches(4.5))

    prs.save(output_filename)
    return output_filename


